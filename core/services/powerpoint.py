import json
import os
from io import BytesIO
from itertools import groupby

from pptx import Presentation
from pptx.util import Inches

from core.services.base import BaseDocumentService


class PowerpointGenerateService(BaseDocumentService):
    def __init__(self):
        self.presentation = None

    def load(self, file) -> None:
        self.presentation = Presentation(file)

    def update(self, data: dict) -> None:
        if not self.presentation:
            raise ValueError("PPTX файл не загружен")

        for key, value in data.items():
            if key.startswith("a") and key[1:].isdigit():
                slide_index = int(key[1:]) - 1
                # Убедиться что слайд есть, иначе создать
                if slide_index >= len(self.presentation.slides):
                    slide = self.presentation.slides.add_slide(
                        self.presentation.slide_layouts[6]
                    )  # Использовать пустой шаблон
                else:
                    slide = self.presentation.slides[slide_index]

                # Создать и заполнить фигуру в зависимости от значения
                if isinstance(value, list):
                    for index, item in enumerate(value):
                        if isinstance(item, list):
                            text = "\n".join(
                                f"{i+1}. {sub_item}" for i, sub_item in enumerate(item)
                            )
                            self.create_shape(slide, text, index + 1)
                        else:
                            self.create_shape(slide, item, index + 1)
                else:
                    self.create_shape(slide, value, index + 1)

    def create_shape(self, slide, text, index):
        """Создать фигуру с текстом на слайде"""
        left = width = height = Inches(1)
        textbox = slide.shapes.add_textbox(left, Inches(index), width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = text

    def save_to_bytes(self) -> BytesIO:
        if not self.presentation:
            raise ValueError("PPTX файл не загружен")
        output = BytesIO()
        self.presentation.save(output)
        output.seek(0)
        return output

    def save_to_file(self, file_path: str) -> None:
        if self.presentation:
            self.presentation.save(file_path)
        else:
            raise ValueError("PPTX файл не загружен")


class PowerpointAnalyzeService:
    def __init__(self):
        self.ppt = None
        self.list_of_notes = []
        self.list_of_notes_address = []

    def load(self, file) -> None:
        self.ppt = Presentation(file)

    # Основная функция работы с презентацией
    def analyze(self):
        # Получение адресов закладок/страниц
        for slide in self.ppt.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        note_address = run.hyperlink.address
                        if note_address != None:
                            self.list_of_notes_address.append(note_address)
                        else:
                            continue
        # Удаление повторяющихся адресов, возникающих из-за пробелов и других символов
        new_list_of_notes_address = [
            el for el, _ in groupby(self.list_of_notes_address)
        ]

        notes = PowerpointAnalyzeService.get_text_from_slides(self.ppt)
        PowerpointAnalyzeService.merging_lists(
            new_list_of_notes_address, notes, self.list_of_notes
        )
        # PowerPointService.save_to_json(list_of_notes)
        return self.list_of_notes

    # Получение текста со слайдов
    def get_text_from_slides(ppt):
        notes = []
        for page, slide in enumerate(ppt.slides):
            temp = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    temp.append(shape.text)
            temp = PowerpointAnalyzeService.check_slash_n(temp)
            notes.append(temp)
        return notes

    # Объединение массива адресов и массива текстов в один
    def merging_lists(list1, list2, list3):
        for i in range(0, max(len(list1), len(list2))):
            if i < len(list1):
                list3.append(list1[i])
            if i < len(list2):
                list3.append(list2[i])
        # print(list3)
        return list3

    # Удаление символов табуляции и перехода на другую строку из массива
    def check_slash_n(list):
        list = [x.replace("\n", " ") for x in list]
        list = [x.replace("\x0b", " ") for x in list]
        list = [x.replace("\r", " ") for x in list]
        list = [x.replace("\t", " ") for x in list]
        return list

    # Сохранение массива в json файл
    def save_to_json(self):
        if os.path.isdir("/output"):
            with open("./output/notes.json", "w", encoding="utf-8") as file:
                json.dump(self.list_of_notes, file, ensure_ascii=False, indent=4)
        else:
            os.makedirs("output", exist_ok=True)
            with open("./output/notes.json", "w", encoding="utf-8") as file:
                json.dump(self.list_of_notes, file, ensure_ascii=False, indent=4)

    # Строка для ввода названия презентации
    def get_pres_name(self):
        print("Enter the name of your presentation: ")
